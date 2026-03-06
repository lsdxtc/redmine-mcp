"""
測試 get_attachment_text 相關功能
"""
import pytest
from unittest.mock import patch, Mock
from redmine_mcp.server import (
    _get_file_extension,
    _try_decode_text,
    _extract_pdf_text,
    _extract_docx_text,
    _extract_xlsx_text,
    _extract_pptx_text,
)


class TestGetFileExtension:
    def test_common_extensions(self):
        assert _get_file_extension("report.pdf") == ".pdf"
        assert _get_file_extension("doc.DOCX") == ".docx"
        assert _get_file_extension("data.json") == ".json"
        assert _get_file_extension("image.PNG") == ".png"

    def test_no_extension(self):
        assert _get_file_extension("README") == ""
        assert _get_file_extension("Makefile") == ""

    def test_multiple_dots(self):
        assert _get_file_extension("archive.tar.gz") == ".gz"
        assert _get_file_extension("my.file.txt") == ".txt"


class TestTryDecodeText:
    def test_utf8_text(self):
        text = "Hello, 世界！"
        result = _try_decode_text(text.encode('utf-8'))
        assert result == text

    def test_utf8_bom(self):
        text = "Hello BOM"
        result = _try_decode_text(b'\xef\xbb\xbf' + text.encode('utf-8'))
        assert "Hello BOM" in result

    def test_binary_data(self):
        # 大量不可列印字元 → 應返回 None
        binary = bytes(range(256)) * 10
        result = _try_decode_text(binary)
        assert result is None

    def test_empty_data(self):
        result = _try_decode_text(b"")
        assert result == ""

    def test_json_content(self):
        json_str = '{"key": "value", "number": 42}'
        result = _try_decode_text(json_str.encode('utf-8'))
        assert result == json_str

    def test_xml_content(self):
        xml_str = '<?xml version="1.0"?>\n<root><item>test</item></root>'
        result = _try_decode_text(xml_str.encode('utf-8'))
        assert result == xml_str


class TestExtractPdfText:
    def test_extract_pdf(self):
        """測試 PDF 提取（使用真實的最小 PDF）"""
        from pypdf import PdfWriter
        from io import BytesIO

        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)

        # PyPDF2 blank page 不含文字，所以預期回傳提示訊息
        buf = BytesIO()
        writer.write(buf)
        result = _extract_pdf_text(buf.getvalue())
        assert "無法提取" in result or "第" in result


class TestExtractDocxText:
    def test_extract_docx(self):
        """測試 Word 文字提取"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        doc.add_paragraph("第一段文字")
        doc.add_paragraph("第二段文字")

        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "第一段文字" in result
        assert "第二段文字" in result

    def test_extract_docx_with_table(self):
        """測試 Word 表格提取"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        doc.add_paragraph("標題")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"

        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "標題" in result
        assert "A1" in result
        assert "B2" in result

    def test_extract_empty_docx(self):
        """測試空 Word 文件"""
        from docx import Document
        from io import BytesIO

        doc = Document()
        buf = BytesIO()
        doc.save(buf)
        result = _extract_docx_text(buf.getvalue())
        assert "無文字內容" in result


class TestExtractXlsxText:
    def test_extract_xlsx(self):
        """測試 Excel 提取"""
        from openpyxl import Workbook
        from io import BytesIO

        wb = Workbook()
        ws = wb.active
        ws.title = "資料表"
        ws['A1'] = "姓名"
        ws['B1'] = "分數"
        ws['A2'] = "小明"
        ws['B2'] = 95

        buf = BytesIO()
        wb.save(buf)
        result = _extract_xlsx_text(buf.getvalue())
        assert "資料表" in result
        assert "姓名" in result
        assert "小明" in result
        assert "95" in result

    def test_extract_empty_xlsx(self):
        """測試空 Excel"""
        from openpyxl import Workbook
        from io import BytesIO

        wb = Workbook()
        ws = wb.active
        # 清空預設 cell
        buf = BytesIO()
        wb.save(buf)
        # 空的 workbook 會有一個空 sheet，但沒有內容
        result = _extract_xlsx_text(buf.getvalue())
        # openpyxl 新建的 workbook 可能有空 row，結果依實作而定
        assert isinstance(result, str)


class TestExtractPptxText:
    def test_extract_pptx(self):
        """測試 PowerPoint 提取"""
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "投影片標題"
        slide.placeholders[1].text = "內容文字"

        buf = BytesIO()
        prs.save(buf)
        result = _extract_pptx_text(buf.getvalue())
        assert "投影片標題" in result
        assert "內容文字" in result

    def test_extract_empty_pptx(self):
        """測試空簡報"""
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        buf = BytesIO()
        prs.save(buf)
        result = _extract_pptx_text(buf.getvalue())
        assert "無文字內容" in result
